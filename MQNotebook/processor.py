# processor.py
import os
import sys
import pandas as pd
import docx
from pptx import Presentation # Direct control over slides
from llama_index.core import SimpleDirectoryReader, VectorStoreIndex, StorageContext, Document
from llama_index.vector_stores.chroma import ChromaVectorStore
from llama_index.core.memory import ChatMemoryBuffer
from llama_index.readers.file import ImageReader
import chromadb
from config import DB_DIR, TEMP_DATA_DIR, Settings

# ==========================================
# 1. HARDCORE READERS (The Heavy Lifters)
# ==========================================

class HardcoreDocxReader:
    """Reads Word files by brute force."""
    def load_data(self, file, extra_info=None):
        doc = docx.Document(file)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_data:
                    full_text.append(" | ".join(row_data))
        
        text = "\n\n".join(full_text)
        return [Document(text=text, extra_info=extra_info or {})]

class HardcorePptxReader:
    """
    Iterates through every shape, text box, table, and speaker note.
    Does not care about layout validation.
    """
    def load_data(self, file, extra_info=None):
        prs = Presentation(file)
        full_text = []
        
        for i, slide in enumerate(prs.slides):
            slide_content = [f"--- Slide {i+1} ---"]
            
            # 1. Extract Text from Shapes (Titles, Body, TextBoxes)
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs).strip()
                        if text:
                            slide_content.append(text)
                
                # 2. Extract Text from Tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text_frame.text.strip() for cell in row.cells if cell.text_frame.text.strip()])
                        if row_text:
                            slide_content.append(f"[Table Row]: {row_text}")
            
            # 3. Extract Speaker Notes (Crucial context often hidden)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                 notes = slide.notes_slide.notes_text_frame.text.strip()
                 if notes:
                     slide_content.append(f"[Speaker Notes]: {notes}")

            # Only add the slide if we found text
            if len(slide_content) > 1: 
                full_text.extend(slide_content)

        final_text = "\n\n".join(full_text)
        
        # Fallback: If PPT is just images, return empty string so we can error handle later
        if not final_text.strip():
            return [] # Empty list tells system "nothing here"
            
        return [Document(text=final_text, extra_info=extra_info or {})]

class HardcoreExcelReader:
    """
    Flattens every sheet into a text format.
    Uses Markdown table format for better LLM readability.
    """
    def load_data(self, file, extra_info=None):
        try:
            # Read all sheets
            dfs = pd.read_excel(file, sheet_name=None)
            full_text = []
            
            for sheet_name, df in dfs.items():
                # clean empty rows/cols
                df = df.dropna(how='all').dropna(axis=1, how='all')
                
                if not df.empty:
                    full_text.append(f"--- Sheet: {sheet_name} ---")
                    # Convert to markdown string (preserves structure for LLM)
                    # We convert to string to ensure it's JSON serializable later if needed
                    full_text.append(df.to_string(index=False))
            
            final_text = "\n\n".join(full_text)
            return [Document(text=final_text, extra_info=extra_info or {})]
            
        except Exception as e:
            print(f"Excel Error: {e}")
            return []

# ==========================================
# 2. PIPELINE LOGIC
# ==========================================

def get_file_extractors():
    return {
        ".docx": HardcoreDocxReader(),
        ".pptx": HardcorePptxReader(), # <-- New Beast
        ".xlsx": HardcoreExcelReader(), # <-- New Beast
        ".jpg": ImageReader(text_type="text"),
        ".png": ImageReader(text_type="text"),
    }

def process_documents(uploaded_files):
    if not uploaded_files:
        return None

    # 1. Create unique temp directory
    if not os.path.exists(TEMP_DATA_DIR):
        os.makedirs(TEMP_DATA_DIR)
        
    for uploaded_file in uploaded_files:
        file_path = os.path.join(TEMP_DATA_DIR, uploaded_file.name)
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

    print(f"📂 Processing in unique folder: {DB_DIR}")

    # 2. Load Data with Hardcore Parsers
    documents = SimpleDirectoryReader(
        input_dir=TEMP_DATA_DIR,
        file_extractor=get_file_extractors(),
        recursive=True
    ).load_data()
    
    # FILTER: Remove empty docs immediately to stop "node missing content" error
    valid_documents = [doc for doc in documents if doc.text and doc.text.strip()]

    if not valid_documents:
        print("⚠️ All uploaded documents appear empty or unreadable.")
        return None

    # 3. Initialize Database
    chroma_client = chromadb.PersistentClient(path=DB_DIR)
    chroma_collection = chroma_client.get_or_create_collection("session_collection")
    vector_store = ChromaVectorStore(chroma_collection=chroma_collection)
    storage_context = StorageContext.from_defaults(vector_store=vector_store)

    # 4. Build Index
    index = VectorStoreIndex.from_documents(
        valid_documents, # Only pass valid docs
        storage_context=storage_context,
        embed_model=Settings.embed_model,
        show_progress=True
    )
    
    return index

def get_chat_engine(index, reranker):
    memory = ChatMemoryBuffer.from_defaults(token_limit=15000)
    
    chat_engine = index.as_chat_engine(
        chat_mode="context",
        memory=memory,
        node_postprocessors=[reranker],
        similarity_top_k=10,
        system_prompt=(
            "You are an MQNotebook Assistant. "
            "Analyze the provided uploaded files (Word, Excel, PPT) carefully. "
            "Data from Excel has been converted to text tables - interpret them row by row. "
            "Data from PPT includes slide numbers and speaker notes. "
            "If the user asks about specific documents, refer to them by name. "
        )
    )
    return chat_engine